# Copyright 2025 Vijaykumar Singh <singhvjd@gmail.com>
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Document processing capabilities for multi-modal AI.

Supports:
- PDF text extraction
- DOCX processing
- PPTX processing
- TXT and Markdown files
- Image extraction from documents
"""

import logging
from pathlib import Path
from typing import TYPE_CHECKING, Any, Dict, List, Optional, Union

if TYPE_CHECKING:
    from victor.multimodal.processor import ProcessingResult

logger = logging.getLogger(__name__)

# Optional dependencies with graceful fallback
try:
    import PyPDF2  # type: ignore[import]

    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False
    logger.warning("PyPDF2 not available. Install with: pip install PyPDF2")

try:
    from docx import Document  # type: ignore[import-not-found]

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available. Install with: pip install python-docx")

try:
    from pptx import Presentation  # type: ignore[import-not-found]

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. Install with: pip install python-pptx")


class DocumentProcessor:
    """Process and analyze documents.

    Example:
        processor = DocumentProcessor()

        # Process PDF
        result = await processor.process(
            document_path="report.pdf",
            query="Summarize key findings"
        )

        # Extract images
        result = await processor.process(
            document_path="presentation.pptx",
            extract_images=True
        )
    """

    supported_formats = [
        ".pdf",
        ".docx",
        ".doc",
        ".pptx",
        ".ppt",
        ".txt",
        ".md",
        ".markdown",
        ".rtf",
        ".csv",
    ]

    def __init__(self, config: Optional[Any] = None):
        """Initialize document processor.

        Args:
            config: ProcessingConfig instance
        """
        self.config = config

    async def process(
        self,
        document_path: Union[str, Path],
        query: Optional[str] = None,
        extract_images: bool = False,
    ) -> "ProcessingResult":
        """Process document and extract text.

        Args:
            document_path: Path to document file
            query: Optional query for document understanding
            extract_images: Whether to extract images from document

        Returns:
            ProcessingResult with document content
        """
        from victor.multimodal.processor import (
            MediaType,
            ProcessingResult,
            ProcessingStatus,
        )

        path = Path(document_path)

        if not path.exists():
            return ProcessingResult(
                status=ProcessingStatus.FAILED,
                content="",
                error=f"Document file not found: {document_path}",
                media_type=MediaType.DOCUMENT,
                source=str(document_path),
            )

        try:
            # Extract text based on format
            ext = path.suffix.lower()

            if ext == ".pdf":
                content, metadata = await self._extract_pdf(path, extract_images)
            elif ext in [".docx", ".doc"]:
                content, metadata = await self._extract_docx(path)
            elif ext in [".pptx", ".ppt"]:
                content, metadata = await self._extract_pptx(path)
            elif ext in [".txt", ".md", ".markdown", ".rtf", ".csv"]:
                content, metadata = await self._extract_text(path)
            else:
                return ProcessingResult(
                    status=ProcessingStatus.UNSUPPORTED,
                    content="",
                    error=f"Unsupported document format: {ext}",
                    media_type=MediaType.DOCUMENT,
                    source=str(document_path),
                )

            # Add query-specific processing if provided
            if query and content:
                content = self._process_with_query(content, query)

            return ProcessingResult(
                status=ProcessingStatus.SUCCESS,
                content=content,
                metadata=metadata,
                confidence=0.95,
                media_type=MediaType.DOCUMENT,
                source=str(document_path),
            )

        except Exception as e:
            logger.error(f"Document processing failed: {e}")
            return ProcessingResult(
                status=ProcessingStatus.FAILED,
                content="",
                error=str(e),
                media_type=MediaType.DOCUMENT,
                source=str(document_path),
            )

    async def _extract_pdf(self, path: Path, extract_images: bool) -> tuple[str, Dict[str, Any]]:
        """Extract text from PDF."""
        if not PYPDF2_AVAILABLE:
            return "", {"error": "PyPDF2 not available"}

        try:
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)

                # Get metadata
                metadata = {
                    "pages": len(reader.pages),
                    "title": reader.metadata.get("/Title", "") if reader.metadata else "",
                    "author": reader.metadata.get("/Author", "") if reader.metadata else "",
                }

                # Extract text from all pages
                text_parts = []
                for i, page in enumerate(reader.pages):
                    text = page.extract_text()
                    text_parts.append(f"--- Page {i + 1} ---\n{text}")

                content = "\n\n".join(text_parts)

                return content, metadata

        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            return "", {"error": str(e)}

    async def _extract_docx(self, path: Path) -> tuple[str, Dict[str, Any]]:
        """Extract text from DOCX."""
        if not DOCX_AVAILABLE:
            return "", {"error": "python-docx not available"}

        try:
            doc = Document(path)

            # Extract paragraphs
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]

            content = "\n\n".join(paragraphs)

            # Extract tables
            tables_data = []
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    table_text.append(" | ".join(row_text))
                tables_data.append("\n".join(table_text))

            metadata = {
                "paragraphs": len(paragraphs),
                "tables": len(tables_data),
                "tables_data": tables_data,
            }

            return content, metadata

        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            return "", {"error": str(e)}

    async def _extract_pptx(self, path: Path) -> tuple[str, Dict[str, Any]]:
        """Extract text from PPTX."""
        if not PPTX_AVAILABLE:
            return "", {"error": "python-pptx not available"}

        try:
            prs = Presentation(path)

            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_content = []

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_content.append(shape.text)

                slides_text.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_content))

            content = "\n\n".join(slides_text)

            metadata = {
                "slides": len(prs.slides),
                "title": (
                    prs.slides[0].shapes.title.text
                    if prs.slides and prs.slides[0].shapes.title
                    else ""
                ),
            }

            return content, metadata

        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            return "", {"error": str(e)}

    async def _extract_text(self, path: Path) -> tuple[str, Dict[str, Any]]:
        """Extract text from plain text files."""
        try:
            # Try different encodings
            encodings = ["utf-8", "latin-1", "cp1252"]

            content = None
            for encoding in encodings:
                try:
                    with open(path, "r", encoding=encoding) as f:
                        content = f.read()
                    break
                except UnicodeDecodeError:
                    continue

            if content is None:
                content = ""

            # Count lines, words, characters
            lines = content.split("\n")
            words = content.split()

            metadata = {
                "lines": len(lines),
                "words": len(words),
                "characters": len(content),
                "encoding": encoding,
            }

            return content, metadata

        except Exception as e:
            logger.error(f"Text extraction failed: {e}")
            return "", {"error": str(e)}

    def _process_with_query(self, content: str, query: str) -> str:
        """Process content with query context.

        This is a simple implementation. For production, use LLM-based analysis.
        """
        # Extract relevant sections based on query keywords
        query_lower = query.lower()
        keywords = query_lower.split()

        # Simple keyword-based highlighting
        lines = content.split("\n")
        relevant_lines = []

        for line in lines:
            line_lower = line.lower()
            if any(keyword in line_lower for keyword in keywords if len(keyword) > 3):
                relevant_lines.append(line)

        if relevant_lines:
            return f"Query: {query}\n\nRelevant Content:\n" + "\n".join(relevant_lines)
        else:
            return f"Query: {query}\n\nFull Content:\n{content}"

    def get_document_summary(self, content: str, max_length: int = 500) -> str:
        """Generate document summary.

        Args:
            content: Document content
            max_length: Maximum summary length

        Returns:
            Document summary
        """
        # Simple extractive summarization
        sentences = content.split(". ")

        # Take first few sentences up to max_length
        summary_sentences = []
        current_length = 0

        for sentence in sentences:
            if current_length + len(sentence) > max_length:
                break
            summary_sentences.append(sentence)
            current_length += len(sentence) + 2  # +2 for ". "

        return ". ".join(summary_sentences)
